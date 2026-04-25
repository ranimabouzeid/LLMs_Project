from pathlib import Path
from typing import List, Dict, Any

from pypdf import PdfReader
from docx import Document
from pptx import Presentation


def load_pdf(file_path: str) -> List[Dict[str, Any]]:
    """
    Loads a PDF and returns a list of pages.
    Each page is represented as a dictionary with text and metadata.
    """
    path = Path(file_path)
    reader = PdfReader(str(path))

    pages = []

    for page_number, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""

        if text.strip():
            pages.append(
                {
                    "text": text.strip(),
                    "source_file": path.name,
                    "page_number": page_number,
                    "file_type": "pdf",
                }
            )

    return pages


def load_docx(file_path: str) -> List[Dict[str, Any]]:
    """
    Loads a DOCX file and returns the extracted text.
    DOCX files do not naturally have page numbers, so page_number is None.
    """
    path = Path(file_path)
    document = Document(str(path))

    paragraphs = []

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if text:
            paragraphs.append(text)

    full_text = "\n".join(paragraphs)

    if not full_text.strip():
        return []

    return [
        {
            "text": full_text.strip(),
            "source_file": path.name,
            "page_number": None,
            "file_type": "docx",
        }
    ]


def load_pptx(file_path: str) -> List[Dict[str, Any]]:
    """
    Loads a PPTX file and returns one text block per slide.
    slide number is stored as page_number for consistency.
    """
    path = Path(file_path)
    presentation = Presentation(str(path))

    slides = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_text_parts = []

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    slide_text_parts.append(text)

        slide_text = "\n".join(slide_text_parts)

        if slide_text.strip():
            slides.append(
                {
                    "text": slide_text.strip(),
                    "source_file": path.name,
                    "page_number": slide_number,
                    "file_type": "pptx",
                }
            )

    return slides


def load_document(file_path: str) -> List[Dict[str, Any]]:
    """
    Detects file type and loads the document.
    Supported formats: PDF, DOCX, PPTX.
    """
    path = Path(file_path)
    suffix = path.suffix.lower()

    if suffix == ".pdf":
        return load_pdf(file_path)

    if suffix == ".docx":
        return load_docx(file_path)

    if suffix == ".pptx":
        return load_pptx(file_path)

    raise ValueError(
        f"Unsupported file type: {suffix}. Supported types are .pdf, .docx, .pptx"
    )


def load_documents_from_folder(folder_path: str) -> List[Dict[str, Any]]:
    """
    Loads all supported files from a folder.
    """
    folder = Path(folder_path)

    if not folder.exists():
        raise FileNotFoundError(f"Folder not found: {folder_path}")

    all_pages = []
    supported_extensions = {".pdf", ".docx", ".pptx"}

    for file_path in folder.iterdir():
        if file_path.is_file() and file_path.suffix.lower() in supported_extensions:
            pages = load_document(str(file_path))
            all_pages.extend(pages)

    return all_pages